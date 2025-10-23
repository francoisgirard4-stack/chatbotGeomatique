# -*- coding: utf-8 -*-
import os, glob
import numpy as np
import streamlit as st
from openai import OpenAI
import tiktoken
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

st.set_page_config(page_title="Assistant Géomatique – Ultralight", page_icon="🗺️", layout="wide")
st.title("🗺️ Assistant pédagogique en géomatique – Ultralight")

api_key = os.getenv("OPENAI_API_KEY", "")
if not api_key:
    st.error("OPENAI_API_KEY manquante. Ajoutez-la dans Secrets (Streamlit Cloud).")
    st.stop()
client = OpenAI(api_key=api_key)

def read_pdf(path):
    try:
        reader = PdfReader(path)
        texts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(t for t in texts if t)
    except Exception as e:
        st.warning(f"PDF non lu ({os.path.basename(path)}): {e}")
        return ""

def read_docx(path):
    try:
        d = DocxDocument(path)
        parts = []
        for p in d.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text.strip())
        for tbl in d.tables:
            for row in tbl.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text))
        return "\n".join(parts)
    except Exception as e:
        st.warning(f"DOCX non lu ({os.path.basename(path)}): {e}")
        return ""

def read_pptx(path):
    try:
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            bits = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    bits.append(shape.text)
            if bits:
                slides.append("\n".join(bits))
        return "\n---\n".join(slides)
    except Exception as e:
        st.warning(f"PPTX non lu ({os.path.basename(path)}): {e}")
        return ""

def cosine_sim(a, b):
    an = np.linalg.norm(a); bn = np.linalg.norm(b)
    if an == 0 or bn == 0:
        return 0.0
    return float(np.dot(a, b) / (an * bn))

def embed_texts(client, texts):
    resp = client.embeddings.create(model="text-embedding-3-small", input=texts)
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]

def count_tokens(s):
    try:
        enc = tiktoken.get_encoding("cl100k_base")
        return len(enc.encode(s))
    except Exception:
        return len(s.split())

base = os.path.join(os.getcwd(), "docs")
if not os.path.exists(base):
    st.error("Le dossier ./docs est introuvable. Créez-le et placez-y vos fichiers (.pdf/.docx/.pptx).")
    st.stop()

paths = glob.glob(os.path.join(base, "*"))
paths = [p for p in paths if p.lower().endswith((".pdf",".docx",".pptx"))]
if not paths:
    st.warning("Aucun fichier détecté dans ./docs."); st.stop()

raw_docs = []
for p in paths:
    if p.lower().endswith(".pdf"): txt = read_pdf(p)
    elif p.lower().endswith(".docx"): txt = read_docx(p)
    elif p.lower().endswith(".pptx"): txt = read_pptx(p)
    else: txt = ""
    if txt.strip(): raw_docs.append({"path": p, "text": txt})

if not raw_docs:
    st.error("Impossible de lire le contenu des documents."); st.stop()

CHUNK, OVERLAP = 1200, 150
chunks = []
for d in raw_docs:
    t = d["text"]; i = 0
    while i < len(t):
        chunk = t[i:i+CHUNK]
        chunks.append({"source": os.path.basename(d["path"]), "text": chunk})
        i += CHUNK - OVERLAP

@st.cache_resource(show_spinner=False)
def build_index(chunks):
    texts = [c["text"] for c in chunks]
    vecs = embed_texts(client, texts)
    return vecs

with st.spinner("Indexation des documents…"):
    vectors = build_index(chunks)

st.sidebar.header("Modes")
admin = st.sidebar.checkbox("Activer le mode administrateur")
quiz = st.sidebar.checkbox("Activer le mode quiz")

if admin:
    st.info("Mode admin : tu peux noter des corrections ci-dessous (simple mémoire de session).")
    cq = st.text_area("Question à corriger")
    ca = st.text_area("Nouvelle réponse / note")
    if st.button("Enregistrer la correction"):
        st.session_state.setdefault("corrections", []).append((cq, ca))
        st.success("Correction enregistrée.")

if quiz:
    st.header("🎯 Quiz")
    q, a = "Projection conforme ?", "Conserve les angles (ex. Mercator)."
    st.subheader(q)
    ans = st.text_input("Ta réponse :")
    if st.button("Vérifier"):
        st.write(f"**Réponse attendue :** {a}")

st.divider()
question = st.text_input("💬 Pose ta question (carto/SIG/GNSS/projections/interpolation) :")
if question:
    with st.spinner("Recherche de contexte…"):
        qvec = embed_texts(client, [question])[0]
        sims = [(cosine_sim(qvec, v), i) for i, v in enumerate(vectors)]
        sims.sort(reverse=True)
        topk = [chunks[i]["text"] for _, i in sims[:4]]
        sources = [chunks[i]["source"] for _, i in sims[:4]]

    context = "\n\n".join(topk)
    while count_tokens(context) > 3500:
        context = context[:len(context)//2]

    system = (
        "Tu es un assistant pédagogique francophone pour deux cours de géomatique "
        "(1re et 2e année). Tu expliques clairement, en citant la source de tes extraits."
    )
    user = f"Question: {question}\n\nContexte (extraits docs):\n{context}\n\nRéponds en français, clair et concis. Cite les noms de fichiers utilisés."

    with st.spinner("Rédaction de la réponse…"):
        chat = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role":"system","content":system},
                {"role":"user","content":user},
            ],
            temperature=0.4,
        )
        answer = chat.choices[0].message.content

    st.markdown("### 🧠 Réponse")
    st.write(answer)
    st.caption("Sources probables : " + ", ".join(sorted(set(sources))[:4]))
